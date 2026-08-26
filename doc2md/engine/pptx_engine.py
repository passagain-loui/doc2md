"""PPTX engine built on python-pptx: slide-by-slide Markdown with speaker notes."""

from __future__ import annotations

from pathlib import Path

from doc2md.core.errors import ConversionError, EngineUnavailableError
from doc2md.core.router import FileKind
from doc2md.engine.base import BaseEngine


class PptxEngine(BaseEngine):
    name = "pptx"
    supported_kinds = (FileKind.PPTX,)
    requires_process_isolation = False

    def convert(self, source: Path, options: dict) -> str:
        self.validate_source(source)
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise EngineUnavailableError(
                "PPTX backend missing: pip install 'doc2md[docs]' (python-pptx)"
            ) from exc

        try:
            presentation = Presentation(str(source))
        except Exception as exc:
            raise ConversionError(f"Corrupted or unreadable PPTX: {source} ({exc})") from exc

        parts: list[str] = [f"# {Path(source).name}", ""]
        try:
            for index, slide in enumerate(presentation.slides, start=1):
                parts.append(self._slide_markdown(index, slide))
                parts.append("")
        except ConversionError:
            raise
        except Exception as exc:
            raise ConversionError(f"PPTX conversion failed: {source} ({exc})") from exc
        return "\n".join(parts).strip() + "\n"

    def _slide_markdown(self, index: int, slide) -> str:
        title = ""
        bullets: list[str] = []
        tables: list[str] = []
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_lines = [
                        " ".join(p.text.split()) for p in shape.text_frame.paragraphs
                    ]
                    text_lines = [t for t in text_lines if t]
                    if self._is_title_shape(slide, shape):
                        title = text_lines[0] if text_lines else ""
                        continue
                    if not text_lines:
                        continue
                    if title == "" and not bullets:
                        title = text_lines.pop(0)
                    for t in text_lines:
                        bullets.append(f"- {self._bullet_text(t)}")
                elif getattr(shape, "has_table", False) and shape.has_table:
                    tables.append(self._table_to_markdown(shape.table))
        except Exception as exc:
            raise ConversionError(f"Failed to parse slide {index}: {exc}") from exc

        heading = f"## Slide {index}: {title}".rstrip(": ")
        out = [heading]
        out.extend(tables)
        out.extend(bullets)
        notes = self._notes_text(slide)
        if notes:
            out.append("")
            out.append(f"> Speaker notes: {notes}")
        return "\n".join(out)

    @staticmethod
    def _is_title_shape(slide, shape) -> bool:
        try:
            if not shape.is_placeholder or not shape.has_text_frame:
                return False
            return shape.placeholder_format.idx == 0
        except Exception:
            return False

    @staticmethod
    def _bullet_text(text: str) -> str:
        return text.replace("\u2022", "").strip()

    def _notes_text(self, slide) -> str:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            return ""
        if notes_slide is None or notes_slide.notes_text_frame is None:
            return ""
        return " ".join(notes_slide.notes_text_frame.text.split())

    def _table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            cells = [" ".join(cell.text.split()).replace("|", "\\|") for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if not rows:
            return ""
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]
        lines = [
            "| " + " | ".join(rows[0]) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        lines.extend("| " + " | ".join(r) + " |" for r in rows[1:])
        return "\n".join(lines)
