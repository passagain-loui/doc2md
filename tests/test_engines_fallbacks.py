import sys

import pytest

from doc2md.core.errors import ConversionError, EngineUnavailableError
from doc2md.engine import get_engine
from doc2md.core.router import FileKind
from tests.helpers import make_fake_pytesseract


@pytest.fixture
def fake_tesseract(monkeypatch):
    monkeypatch.setitem(sys.modules, "pytesseract", make_fake_pytesseract())
    monkeypatch.setattr(
        "doc2md.engine.ocr_engine.shutil.which", lambda name: "C:/fake/tesseract.exe"
    )


def test_pdf_engine_unavailable(monkeypatch, tmp_path):
    monkeypatch.setitem(sys.modules, "pymupdf", None)
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4")
    with pytest.raises(EngineUnavailableError):
        get_engine(FileKind.PDF).convert(p, {})


def test_docx_engine_unavailable_and_corrupted(monkeypatch, tmp_path):
    p = tmp_path / "garbage.docx"
    p.write_bytes(b"not a zip at all")
    with pytest.raises(EngineUnavailableError):
        monkeypatch.setitem(sys.modules, "docx", None)
        get_engine(FileKind.DOCX).convert(p, {})
    monkeypatch.undo()
    with pytest.raises(ConversionError):
        get_engine(FileKind.DOCX).convert(p, {})


def test_xlsx_engine_unavailable_legacy_and_corrupted(monkeypatch, tmp_path):
    legacy = tmp_path / "old.xls"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0legacy binary")
    with pytest.raises(ConversionError, match="Legacy"):
        get_engine(FileKind.XLSX).convert(legacy, {})

    x = tmp_path / "book.xlsx"
    x.write_bytes(b"PK\x03\x04broken-zip-body")
    monkeypatch.setitem(sys.modules, "openpyxl", None)
    with pytest.raises(EngineUnavailableError):
        get_engine(FileKind.XLSX).convert(x, {})
    monkeypatch.undo()
    with pytest.raises(ConversionError):
        get_engine(FileKind.XLSX).convert(x, {})


def test_pptx_engine_unavailable_and_corrupted(monkeypatch, tmp_path):
    p = tmp_path / "deck.pptx"
    p.write_bytes(b"PK\x03\x04not-a-presentation")
    monkeypatch.setitem(sys.modules, "pptx", None)
    with pytest.raises(EngineUnavailableError):
        get_engine(FileKind.PPTX).convert(p, {})
    monkeypatch.undo()
    with pytest.raises(ConversionError):
        get_engine(FileKind.PPTX).convert(p, {})


def test_ocr_engine_pil_unavailable(monkeypatch, tmp_path):
    from PIL import Image

    p = tmp_path / "a.png"
    Image.new("RGB", (10, 10)).save(p)
    monkeypatch.setitem(sys.modules, "PIL", None)
    with pytest.raises(EngineUnavailableError):
        get_engine(FileKind.IMAGE).convert(p, {})


def test_ocr_engine_oserror_on_open(monkeypatch, tmp_path):
    PIL_Image = pytest.importorskip("PIL.Image")
    p = tmp_path / "ok.png"
    from PIL import Image as RealImage

    RealImage.new("RGB", (10, 10)).save(p)

    def raise_oserror(*args, **kwargs):
        raise OSError("disk failure mid-read")

    monkeypatch.setattr(PIL_Image, "open", raise_oserror)
    with pytest.raises(ConversionError, match="Unreadable image"):
        get_engine(FileKind.IMAGE).convert(p, {})


def test_ocr_success_with_tesseract_stub(fake_tesseract, tmp_path):
    from PIL import Image

    p = tmp_path / "scan.png"
    Image.new("RGB", (80, 40), "white").save(p)
    out = get_engine(FileKind.IMAGE).convert(p, {"ocr_lang": "eng"})
    assert "## Extracted text" in out
    assert "OCR ENGINE TEXT" in out


def test_ocr_empty_text_note(fake_tesseract, monkeypatch, tmp_path):
    from PIL import Image

    monkeypatch.setitem(sys.modules, "pytesseract", make_fake_pytesseract(text="   \n"))
    p = tmp_path / "blank.png"
    Image.new("RGB", (30, 30), "black").save(p)
    out = get_engine(FileKind.IMAGE).convert(p, {})
    assert "no readable text was found" in out


def test_ocr_tesseract_error_wrapped(monkeypatch, tmp_path):
    from PIL import Image

    failing = make_fake_pytesseract()

    def raise_err(path, lang="eng"):
        raise failing.TesseractError(-11, "engine exploded")

    failing.image_to_string = raise_err
    monkeypatch.setitem(sys.modules, "pytesseract", failing)
    monkeypatch.setattr(
        "doc2md.engine.ocr_engine.shutil.which", lambda name: "C:/fake/tesseract.exe"
    )
    p = tmp_path / "e.png"
    Image.new("RGB", (20, 20), "white").save(p)
    with pytest.raises(ConversionError, match="Tesseract OCR failed"):
        get_engine(FileKind.IMAGE).convert(p, {})


def test_scanned_pdf_hint_without_ocr_backend(tmp_path, monkeypatch):
    pymupdf = pytest.importorskip("pymupdf")
    p = tmp_path / "scanned.pdf"
    doc = pymupdf.open()
    for _ in range(2):
        doc.new_page()
    doc.save(str(p))
    doc.close()

    real_which = __import__("shutil").which
    import doc2md.engine.pdf_engine as pe

    monkeypatch.setattr(pe.shutil, "which", lambda name: None)
    out = get_engine(FileKind.PDF).convert(p, {"pdf_ocr_fallback": True})
    assert "No extractable text found on 2 page(s)" in out
    assert "OCR" in out


def test_scanned_pdf_ocr_fallback_renders_pages(tmp_path, monkeypatch, fake_tesseract):
    pymupdf = pytest.importorskip("pymupdf")
    p = tmp_path / "scanned_ocr.pdf"
    doc = pymupdf.open()
    doc.new_page()
    doc.save(str(p))
    doc.close()

    out = get_engine(FileKind.PDF).convert(p, {})
    assert "## Page 1 (OCR)" in out
    assert "OCR ENGINE TEXT" in out


def test_pdf_ocr_fallback_disabled_yields_hint(tmp_path, monkeypatch):
    pymupdf = pytest.importorskip("pymupdf")
    p = tmp_path / "noff.pdf"
    doc = pymupdf.open()
    doc.new_page()
    doc.save(str(p))
    doc.close()
    out = get_engine(FileKind.PDF).convert(p, {"pdf_ocr_fallback": False})
    assert "(OCR)" not in out
    assert "No extractable text" in out


def test_pdf_page_text_failure_wrapped(tmp_path, simple_pdf, monkeypatch):
    pytest.importorskip("pymupdf")
    import pymupdf

    def boom(self, *args, **kwargs):
        raise RuntimeError("text layer exploded")

    monkeypatch.setattr(pymupdf.Page, "get_text", boom)
    with pytest.raises(ConversionError, match="PDF conversion failed"):
        get_engine(FileKind.PDF).convert(simple_pdf, {})


def test_docx_style_variants(tmp_path):
    import docx

    d = docx.Document()
    d.add_heading("Doc Title", level=0)
    d.add_paragraph("numbered item one", style="List Number")
    d.add_paragraph("quoted wisdom", style="Quote")
    t = d.add_table(rows=1, cols=1)
    t.rows[0].cells[0].text = ""
    p = tmp_path / "styles.docx"
    d.save(str(p))
    out = get_engine(FileKind.DOCX).convert(p, {})
    assert "# Doc Title" in out
    assert "1. numbered item one" in out
    assert "> quoted wisdom" in out


def test_docx_table_to_markdown_static_padding():
    from doc2md.engine.docx_engine import DocxEngine

    class Cell:
        def __init__(self, text):
            self._text = text

        @property
        def text(self):
            return self._text

    class Row:
        def __init__(self, cells):
            self.cells = cells

    class Table:
        rows = [Row([Cell("h1")]), Row([Cell("a"), Cell("b")])]

    md = DocxEngine._table_to_markdown(Table())
    assert "| h1 |  |" in md
    assert "| a | b |" in md


def test_excel_empty_sheet_ragged_rows_none_cells(tmp_path):
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.append(["only_header"])
    ws1.append([None, "second"])
    ws1.append(["v", "w"])
    wb.create_sheet("EmptySheet")
    wb.create_sheet("BlankCells")
    wb["BlankCells"]["A1"] = None
    wb["BlankCells"]["B1"] = "kept"
    p = tmp_path / "edge.xlsx"
    wb.save(str(p))

    out = get_engine(FileKind.XLSX).convert(p, {})
    assert "col2" in out
    assert "_(empty sheet)_" in out
    assert "| col1 | kept |" in out


def test_excel_csv_custom_encoding_option(tmp_path):
    p = tmp_path / "thai.csv"
    payload = "หัวตาราง,ค่า\nแถวหนึ่ง,สอง"
    p.write_bytes(payload.encode("cp874"))
    out = get_engine(FileKind.XLSX).convert(
        p, {"encodings": ("cp874", "utf-8"), "max_rows": 100}
    )
    assert "หัวตาราง" in out and "สอง" in out


def test_code_engine_fence_collision(tmp_path):
    p = tmp_path / "md_like.txt"
    p.write_text("# has\n```python\nprint('fenced')\n```\ninside", encoding="utf-8")
    out = get_engine(FileKind.CODE).convert(p, {})
    assert out.startswith("# md_like.txt\n\n````text\n")
    assert out.rstrip().endswith("````")
    assert "print('fenced')" in out


def test_pptx_full_deck_features(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Quarterly Review"
    s1.placeholders[1].text = "Subtitle line here"
    s1.notes_slide.notes_text_frame.text = "Remember the KPIs please"
    grid = s1.shapes.add_table(2, 2, Inches(0.5), Inches(4), Inches(4), Inches(1)).table
    grid.cell(0, 0).text = "Metric|X"
    grid.cell(0, 1).text = "Value"
    grid.cell(1, 0).text = "uptime"
    grid.cell(1, 1).text = "99.9%"

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    box = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
    box.text_frame.text = "Implicit Title From First Bullet"
    box.text_frame.add_paragraph().text = "Second bullet point"

    empty = prs.slides.add_slide(prs.slide_layouts[6])

    p = tmp_path / "full.pptx"
    prs.save(str(p))
    out = get_engine(FileKind.PPTX).convert(p, {})

    assert "## Slide 1: Quarterly Review" in out
    assert "- Subtitle line here" in out
    assert "> Speaker notes: Remember the KPIs please" in out
    assert "| Metric\\|X | Value |" in out
    assert "| uptime | 99.9% |" in out
    assert "## Slide 2: Implicit Title From First Bullet" in out
    assert "- Second bullet point" in out
    assert "## Slide 3\n" in out


def test_web_tag_rendering_matrix():
    from bs4 import BeautifulSoup

    from doc2md.engine.web_engine import WebEngine

    we = WebEngine()

    def render(html_snippet):
        soup = BeautifulSoup(html_snippet, "html.parser")
        return we._render_node(soup.find(True) or soup.contents[0])

    quoted = we._render_node(
        BeautifulSoup("<blockquote>wisdom <b>bold</b></blockquote>", "html.parser").blockquote
    )
    assert quoted.splitlines()[0].startswith("> wisdom")
    assert "**bold**" in quoted

    assert we._render_node(BeautifulSoup("<hr/>", "html.parser").hr) == "---"
    ol = we._render_node(
        BeautifulSoup("<ol><li>one</li><li>two</li></ol>", "html.parser").ol
    )
    assert "1. one" in ol and "2. two" in ol

    em = we._render_node(BeautifulSoup("<p><em>soft</em> <code>x=1</code></p>", "html.parser").p)
    assert "*soft*" in em and "`x=1`" in em

    assert we._render_node(BeautifulSoup('<img alt="no src"/>', "html.parser").img) == ""
    assert we._render_node(BeautifulSoup('<a>bare link</a>', "html.parser").a) == "bare link"

    heading_empty = we._render_node(BeautifulSoup("<h3>   </h3>", "html.parser").h3)
    assert heading_empty == ""

    passthrough = we._render_node(
        BeautifulSoup("<span><font>styled</font></span>", "html.parser").span
    )
    assert "styled" in passthrough

    pre = we._render_node(
        BeautifulSoup("<pre>line1\nline2</pre>", "html.parser").pre
    )
    assert "```" in pre and "line2" in pre

    table_md = we._render_node(
        BeautifulSoup("<table><tr><th>H</th></tr><tr><td>r1c1</td><td>r1c2</td></tr></table>",
                      "html.parser").table
    )
    assert "| H |  |" in table_md
    assert "| r1c1 | r1c2 |" in table_md

    comment_soup = BeautifulSoup("<div><!--hidden note-->visible</div>", "html.parser")
    rendered = we._render_children(comment_soup.div)
    assert "hidden note" not in rendered
    assert "visible" in rendered


def test_web_html_parser_fallback_when_lxml_fails(monkeypatch, tmp_path):
    import bs4

    original = bs4.BeautifulSoup

    def picky(markup, parser=None):
        if parser == "lxml":
            raise RuntimeError("lxml unavailable")
        return original(markup, parser or "html.parser")

    monkeypatch.setattr(bs4, "BeautifulSoup", picky)
    p = tmp_path / "fallback.html"
    p.write_text("<html><head><title>Fallback</title></head><body><p>parsed anyway</p></body></html>",
                 encoding="utf-8")
    out = get_engine(FileKind.HTML).convert(p, {})
    assert "# Fallback" in out and "parsed anyway" in out


def test_web_bs4_missing_raises(monkeypatch, tmp_path):
    monkeypatch.setitem(sys.modules, "bs4", None)
    p = tmp_path / "page.html"
    p.write_text("<p>hi</p>", encoding="utf-8")
    with pytest.raises(EngineUnavailableError):
        get_engine(FileKind.HTML).convert(p, {})


def test_web_html_to_text_regex_fallback_without_bs4(monkeypatch):
    from doc2md.engine.web_engine import WebEngine

    monkeypatch.setitem(sys.modules, "bs4", None)
    out = WebEngine._html_to_text("<p>plain <b>strip</b></p>")
    assert "<" not in out and "strip" in out


def test_eml_with_attachments(tmp_path):
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["From"] = "files@example.com"
    msg["Subject"] = "With attachment"
    msg.set_content("see attached")
    msg.add_attachment(b"%PDF-1.4 fake", maintype="application", subtype="pdf", filename="report.pdf")
    p = tmp_path / "attach.eml"
    p.write_bytes(bytes(msg))
    out = get_engine(FileKind.EML).convert(p, {})
    assert "## Attachments" in out
    assert "- report.pdf" in out


def test_eml_html_only_body(tmp_path):
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["From"] = "rich@example.com"
    msg["Subject"] = "HTML body only"
    msg.set_content("<html><body><p>Rich content lives here</p></body></html>", subtype="html")
    p = tmp_path / "html.eml"
    p.write_bytes(bytes(msg))
    out = get_engine(FileKind.EML).convert(p, {})
    assert "Rich content lives here" in out


def test_eml_bogus_charset_falls_back_to_detector(tmp_path):
    thai_body = "ชุดอักขระเสียหายยังถอดรหัสได้จากข้อมูลที่เพียงพอ ".encode("cp874") * 6
    raw = (
        b"From: broken@example.co.th\r\nTo: me@example.com\r\nSubject: bad charset\r\n"
        b"MIME-Version: 1.0\r\nContent-Type: text/plain; charset=\"x-not-a-real-codepage\"\r\n\r\n"
        + thai_body
    )
    p = tmp_path / "badcharset.eml"
    p.write_bytes(raw)
    out = get_engine(FileKind.EML).convert(p, {})
    assert "ถอดรหัสได้" in out.replace("\n", "")


def test_eml_parser_failure_wrapped(tmp_path, monkeypatch):
    import email

    def explode(*args, **kwargs):
        raise ValueError("malformed beyond parsing")

    monkeypatch.setattr(email, "message_from_bytes", explode)
    p = tmp_path / "junk.eml"
    p.write_bytes(b"\x00\xff\xfe garbage")
    with pytest.raises(ConversionError, match="Unreadable EML"):
        get_engine(FileKind.EML).convert(p, {})


def test_eml_header_exception_swallowed():
    from doc2md.engine.web_engine import WebEngine

    class Boom:
        def get(self, *args, **kwargs):
            raise RuntimeError("header access failed")

    assert WebEngine._header(Boom(), "From") == ""


def test_looks_like_eml_variants(tmp_path):
    from doc2md.engine.web_engine import _looks_like_eml

    yes = tmp_path / "raw.dat"
    yes.write_text("From: a@b.c\r\nSubject: hi\r\n\r\nbody", encoding="utf-8")
    no = tmp_path / "other.dat"
    no.write_bytes(b"\xff\xfe\x00\x01binary")
    assert _looks_like_eml(yes)
    assert not _looks_like_eml(no)
