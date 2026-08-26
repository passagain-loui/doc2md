import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

import pytest
import pymupdf


@pytest.fixture
def simple_pdf(tmp_path):
    path = tmp_path / "simple.pdf"
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello doc2md from PDF")
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def simple_docx(tmp_path):
    import docx

    path = tmp_path / "simple.docx"
    d = docx.Document()
    d.add_heading("Report Title", level=1)
    d.add_paragraph("Intro paragraph with details.")
    d.add_heading("Section", level=2)
    d.add_paragraph("First item", style="List Bullet")
    table = d.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "42"
    d.save(str(path))
    return path


@pytest.fixture
def simple_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "simple.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Deck Title"
    bullets = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(4), Inches(2))
    tf = bullets.text_frame
    tf.text = "Point one"
    tf.add_paragraph().text = "Point two"
    prs.save(str(path))
    return path


@pytest.fixture
def simple_xlsx(tmp_path):
    import openpyxl

    path = tmp_path / "simple.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["col_a", "col_b"])
    ws.append([1, "x|y"])
    ws.append([2, None])
    wb.save(str(path))
    return path
